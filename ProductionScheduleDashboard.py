import streamlit as st
import calendar
import json
import time
from pathlib import Path
from datetime import date, datetime, timedelta
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import io
import re
import holidays

# Optional deps for PPT/Excel — handled later
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
try:
    from pptx.enum.text import MSO_AUTO_SIZE  # may fail on some versions
except Exception:
    MSO_AUTO_SIZE = None

# =======================
# CONFIG & CONSTANTS
# =======================
st.set_page_config(page_title="Production Schedule Dashboard", layout="wide")

COLOR_AC225_RUN_EVG = "#A2EBCD"
COLOR_IN111_RUN_EVG = "#6EC6A8"
COLOR_AC225_RUN_SRX = "#F1E183"
COLOR_IN111_RUN_SRX = "#EC712A"
COLOR_CARDINAL_TPI_NIOWAVE = "#0ABB21"
COLOR_NMCTG = "#D0B9E6"
COLOR_PLACEHOLDER = "#F1E429"
COLOR_SHUTDOWN = "#F5253A"
COLOR_CONFIRMED = "#FF0000"
COLOR_PV = "#3CD63C"
COLOR_SRX = "#3ACCC0"
COLOR_PERCEPTIVE = "#9A6DC1"
COLOR_BWXT = "#3D6E34"
COLOR_MD = "#CC3366"
COLOR_FALLBACK = "#87CEEB"
COLOR_WEEKEND = "#FEF3C7"
COLOR_US_HOLIDAY = "#FF6F3C"
COLOR_CANCELLED = "#E5E7EB"

DASHBOARD_NAME = "production_schedule"
FILENAME = f"{DASHBOARD_NAME}.json"
CONFIG_FILE = Path.home() / ".production_schedule_config.json"
DEFAULT_DIR = Path.home() / "Schedules"
RERUN_FLAG = "__do_rerun__"

# =======================
# SESSION INIT
# =======================
ss = st.session_state
ss.setdefault("current_month", date.today().month)
ss.setdefault("current_year", date.today().year)
ss.setdefault("week_action_rows", {})
ss.setdefault("entries", {})
ss.setdefault("custom_closures", [])
ss.setdefault("__autosave_ok__", False)
ss.setdefault("__autosave_error__", "")
ss.setdefault("show_legend", False)
ss.setdefault("__disk_mtime__", None)
ss.setdefault("suppressed_us_holidays", [])
if "custom_legend_entries" not in ss:
    ss.custom_legend_entries = []
ss.setdefault("__pending_entries__", None)
ss.setdefault("__pending_meta__", None)
ss.setdefault("__pending_week_action_rows__", None)


# =======================
# UTIL & PERSISTENCE
# =======================
def rerun():
    st.rerun()

def _sanitize_dir(raw: str) -> Path:
    s = (raw or "").strip().strip('"').strip("'")
    return Path(s).expanduser()

def load_latest_dir() -> Path:
    if "latest_dir" in ss:
        return Path(ss.latest_dir)
    if CONFIG_FILE.exists():
        try:
            cfg = json.loads(CONFIG_FILE.read_text(encoding="utf-8"))
            last = cfg.get("latest_dir")
            if last:
                ss.latest_dir = last
                return Path(last)
        except Exception:
            pass
    ss.latest_dir = str(DEFAULT_DIR)
    return DEFAULT_DIR

def save_latest_dir(dir_str: str) -> None:
    ss.latest_dir = dir_str
    try:
        CONFIG_FILE.write_text(json.dumps({"latest_dir": dir_str}, indent=2), encoding="utf-8")
    except Exception as e:
        st.warning(f"Couldn't persist latest directory: {e}")

def date_key(y: int, m: int, d: int, row_idx: int) -> str:
    return f"{date(y, m, d).isoformat()}_{row_idx}"

def _save_payload() -> dict:
    return {
        "meta": {"year": ss.current_year, "month": ss.current_month},
        "entries": ss.get("entries", {}),
        "week_action_rows": ss.get("week_action_rows", {}),
        "custom_legend_entries": ss.get("custom_legend_entries", []),
        "suppressed_us_holidays": ss.get("suppressed_us_holidays", []),
    }

def _get_json_path() -> Path:
    latest_dir = _sanitize_dir(str(load_latest_dir()))
    return latest_dir / FILENAME

def _stat_mtime(p: Path):
    try:
        return p.stat().st_mtime if p.exists() else None
    except Exception:
        return None

def _autosave_now() -> Path:
    try:
        payload = _save_payload()
        fp = _get_json_path()
        fp.parent.mkdir(parents=True, exist_ok=True)
        fp.write_text(json.dumps(payload, indent=2, default=str), encoding="utf-8")
        ss["__autosave_ok__"] = True
        ss["__autosave_error__"] = ""
        ss["__disk_mtime__"] = _stat_mtime(fp)
        st.toast("💾 Auto-saved to disk", icon="✅")
        return fp
    except Exception as e:
        ss["__autosave_ok__"] = False
        ss["__autosave_error__"] = str(e)
        raise

def _try_load_from(path: Path):
    try:
        if not path.exists():
            return None, None, None, None
        data = json.loads(path.read_text(encoding="utf-8")) or {}
        if isinstance(data, dict) and "entries" in data:
            entries = data.get("entries", {}) or {}
            meta = data.get("meta") or {}
            week_action_rows = data.get("week_action_rows", {}) or {}
            return entries, meta, week_action_rows, data
    except Exception as e:
        st.error(f"Load error: {e}")
    return None, None, None, None

def _apply_meta_to_calendar(meta: dict) -> bool:
    try:
        yr = int(meta.get("year"))
        mo = int(meta.get("month"))
        changed = (yr != ss.current_year) or (mo != ss.current_month)
        ss.current_year = yr
        ss.current_month = mo
        return changed
    except Exception:
        return False

def _preload_widgets_from_entries():
    for k, v in ss.entries.items():
        ss[f"cell_widget_{k}"] = v.get("text", "")

def _sync_widgets_with_entries():
    for k, v in ss.entries.items():
        wk = f"cell_widget_{k}"
        text_val = v.get("text", "")
        if wk not in ss or ss[wk] != text_val:
            ss[wk] = text_val

def _load_from_disk_into_state() -> bool:
    p = _get_json_path()
    entries, meta, week_action_rows, full_data = _try_load_from(p)
    if entries is None:
        return False

    # --- 🔧 Normalize entries: upgrade legacy strings to {"text": ..., "cancelled": False} ---
    normalized_entries = {}
    for k, v in (entries or {}).items():
        if isinstance(v, dict):
            # Ensure required keys exist
            if "text" not in v:
                v["text"] = ""
            if "cancelled" not in v:
                v["cancelled"] = False
            normalized_entries[k] = v
        elif isinstance(v, str):
            # Convert old string format to new dict format
            normalized_entries[k] = {"text": v.strip(), "cancelled": False}
        else:
            # Fallback
            normalized_entries[k] = {"text": str(v) if v is not None else "", "cancelled": False}

    ss.entries = normalized_entries
    # ---

    if week_action_rows is not None:
        ss.week_action_rows = week_action_rows or {}

    if full_data and "custom_legend_entries" in full_data:
        ss.custom_legend_entries = full_data["custom_legend_entries"] or []

    if full_data and "suppressed_us_holidays" in full_data:
        ss.suppressed_us_holidays = full_data["suppressed_us_holidays"]

    changed = _apply_meta_to_calendar(meta or {})
    _preload_widgets_from_entries()  # Now safe to call
    ss["__disk_mtime__"] = _stat_mtime(p)
    
    if changed:
        ss[RERUN_FLAG] = True

    return True

def _disk_watchdog():
    p = _get_json_path()
    m = _stat_mtime(p)
    if m is None:
        return
    last = ss.get("__disk_mtime__")
    if last is None:
        ss["__disk_mtime__"] = m
        return
    if m > last:
        if _load_from_disk_into_state():
            ss["__disk_mtime__"] = m
            ss[RERUN_FLAG] = True

# =======================
# COLOR / LEGEND
# =======================
def get_color(entry: dict) -> str:
    if not entry or not isinstance(entry, dict):
        return "white"

    text = str(entry.get("text", "")).strip()
    cancelled = entry.get("cancelled", False)

    if not text:
        return "white"

    if cancelled:
        return COLOR_CANCELLED

    lower = text.lower()

    for item in ss.get("custom_legend_entries", []):
        if item["label"].strip().lower() in lower:
            return item["color"]

    if lower == "weekend":
        return COLOR_WEEKEND

    us_holidays = holidays.US(years=ss.current_year)
    if lower in [h.lower() for h in us_holidays.values()]:
        return COLOR_US_HOLIDAY

    for closure in ss.get("custom_closures", []):
        if closure["name"].strip().lower() == lower:
            return COLOR_US_HOLIDAY

    if "shutdown" in lower:
        return COLOR_SHUTDOWN
    if re.search(r"in111.*run.*srx", lower):
        return COLOR_IN111_RUN_SRX
    if re.search(r"ac225.*run.*srx", lower):
        return COLOR_AC225_RUN_SRX
    if re.search(r"in111.*run.*evg", lower):
        return COLOR_IN111_RUN_EVG
    if re.search(r"ac225.*run.*evg", lower):
        return COLOR_AC225_RUN_EVG
    if lower.startswith(("cardinal", "tpi", "niowave")):
        return COLOR_CARDINAL_TPI_NIOWAVE
    if "nmctg" in lower:
        return COLOR_NMCTG
    if re.match(r"^\d{5}-p\d", lower):
        return COLOR_PLACEHOLDER
    if re.match(r"^\d{5}-\d{3}", lower):
        if re.search(r"md[123]$", lower, re.IGNORECASE):
            return COLOR_MD
        return COLOR_CONFIRMED
    if lower.startswith("pv") and "srx" in lower:
        return COLOR_PV
    if lower == "srx maintenance":
        return COLOR_SRX
    if "perceptive" in lower:
        return COLOR_PERCEPTIVE
    if lower == "bwxt order":
        return COLOR_BWXT

    return COLOR_FALLBACK

def is_light_color(hex_color: str) -> bool:
    """Determine if a hex color is light (use black text) or dark (use white text)."""
    if not hex_color or hex_color == "white":
        return True
    if hex_color == "transparent" or hex_color == "none":
        return False
    # Remove '#' if present
    hex_color = hex_color.lstrip('#')
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        # Relative luminance formula (standard for WCAG)
        luminance = (0.299 * r + 0.587 * g + 0.114 * b)
        return luminance > 140  # Threshold: tweak if needed (140 is good for readability)
    except Exception:
        return True  # Default to black text on error

# =======================
# CLINICAL DOSING HELPERS
# =======================
CONFIRMED_PATIENT_CAPTURE = re.compile(r"^\s*(\d{5}-\d{3})\b.*", re.IGNORECASE)

def _extract_patient_code(txt: str):
    if not txt:
        return None
    m = CONFIRMED_PATIENT_CAPTURE.match(str(txt).strip())
    return m.group(1) if m else None

def _is_maintenance(txt: str) -> bool:
    return bool(re.search(r"\bmd[123]\b", (txt or ""), re.IGNORECASE))

def _parse_date_from_dkey(dkey: str):
    try:
        return datetime.fromisoformat(dkey.split("_", 1)[0]).date()
    except Exception:
        return None

def _month_weeks_ext(y: int, m: int):
    cal_raw = calendar.monthcalendar(y, m)
    valid_weeks = [w for w in cal_raw if any(d != 0 for d in w)]
    weeks = []
    for week in valid_weeks:
        ref_day = next((d for d in week if d != 0), None)
        if ref_day is None:
            weeks.append([None] * 7); continue
        ref_date = date(y, m, ref_day)
        start = ref_date - timedelta(days=ref_date.weekday())
        weeks.append([start + timedelta(days=i) for i in range(7)])
    return weeks

def _week_index_for(y: int, m: int, target: date) -> int:
    weeks = _month_weeks_ext(y, m)
    for idx, wk in enumerate(weeks):
        for d in wk:
            if d == target:
                return idx
    return 0

def _first_empty_row_for_date(target: date) -> int:
    y, m, d = target.year, target.month, target.day
    w_idx = _week_index_for(y, m, target)
    key = f"{y}-{m}_{w_idx}"
    num_rows = ss.week_action_rows.get(key, 1)

    for r in range(num_rows):
        dk = date_key(y, m, d, r)
        raw_entry = ss.entries.get(dk)

        # Normalize entry
        if raw_entry is None:
            return r  # truly missing → empty
        if isinstance(raw_entry, str):
            text = raw_entry.strip()
        else:
            text = raw_entry.get("text", "").strip()

        # Consider these as "empty"
        if not text or text.lower() in ["weekend", "placeholder"]:
            return r

    # All current rows are taken → return next row (will trigger row expansion)
    return num_rows

def _entry_exists_on_date(target: date, predicate) -> bool:
    y, m = target.year, target.month
    w_idx = _week_index_for(y, m, target)
    key = f"{y}-{m}_{w_idx}"
    rows = ss.week_action_rows.get(key, 1)
    rows_to_check = rows + 3
    for r in range(rows_to_check):
        dk = date_key(y, m, target.day, r)
        entry = ss.entries.get(dk)
        text_val = entry.get("text", "") if isinstance(entry, dict) else str(entry)
        if text_val and predicate(text_val):
            return True
    return False

def _add_entry_if_absent(target: date, text: str) -> None:
    text_norm = text.strip().lower()
    
    # Don't add if already exists
    if _entry_exists_on_date(target, lambda s: s.strip().lower() == text_norm):
        return
    
    # Find the first empty row
    r = _first_empty_row_for_date(target)
    dk = date_key(target.year, target.month, target.day, r)
    ss.entries[dk] = text.strip()
    
    # Sync widget
    widget_key = f"cell_widget_{dk}"
    if widget_key not in ss:
        ss[widget_key] = text.strip()

def _calc_maintenance_dates(initial_dt: date, n_maint: int = 3, interval_weeks: int = 6):
    dates = []
    cur = initial_dt
    for _ in range(n_maint):
        cur = cur + timedelta(weeks=interval_weeks)
        dates.append(cur)
    return dates

def _cycle_already_scheduled(patient_code: str, initial_dt: date, interval_weeks: int = 6) -> bool:
    md1_date = initial_dt + timedelta(weeks=interval_weeks)
    target_prefix = (patient_code or "").strip().lower()
    return _entry_exists_on_date(
        md1_date,
        lambda s: s.strip().lower().startswith(target_prefix) and re.search(r"\bmd1\b", s, re.IGNORECASE)
    )

def _schedule_patient_cycle(patient_code: str, initial_dt: date, n_maint: int = 3, interval_weeks: int = 6, base_text: str = None):
    if not patient_code:
        return
    # Only schedule maintenance doses if "AC" is in the base text
    if not base_text or 'ac' not in str(base_text).lower():
        return
    if _cycle_already_scheduled(patient_code, initial_dt, interval_weeks=interval_weeks):
        return
    base = (base_text or "").strip() or patient_code
    base_clean = re.sub(r'\b(?:MD[123]|Initial(?:\s*Dose)?)\b', '', base, flags=re.IGNORECASE)
    base_clean = re.sub(r'\s*[-–—]\s*$', '', base_clean).strip()
    if not re.search(rf'\b{re.escape(patient_code)}\b', base_clean):
        base_clean = f"{patient_code} - {base_clean}" if base_clean else patient_code
    for i, dtm in enumerate(_calc_maintenance_dates(initial_dt, n_maint, interval_weeks), start=1):
        _add_entry_if_absent(dtm, f"{base_clean} MD{i}")
    _autosave_now()

def _delete_maintenance_doses(patient_code: str):
    """Delete MD1, MD2, MD3 entries for the given patient code."""
    if not patient_code:
        return
    patient_code_lower = patient_code.strip().lower()
    keys_to_remove = []

    # Find all entries that match patient_code and are MD1/MD2/MD3
    for k, entry in ss.entries.items():
        if isinstance(entry, dict):
            text = entry.get("text", "")
        else:
            text = str(entry)
        text_lower = text.lower()
        if (patient_code_lower in text_lower) and re.search(r"\bmd[123]\b", text_lower, re.IGNORECASE):
            keys_to_remove.append(k)

    # Remove entries and their widgets
    for k in keys_to_remove:
        ss.entries.pop(k, None)
        widget_key = f"cell_widget_{k}"
        if widget_key in ss:
            del ss[widget_key]

def _find_maintenance_doses(patient_code: str):
    """Find all MD1/MD2/MD3 entries for the given patient code."""
    if not patient_code:
        return []
    patient_code_lower = patient_code.strip().lower()
    matches = []
    for k, entry in ss.entries.items():
        if isinstance(entry, dict):
            text = entry.get("text", "")
        else:
            text = str(entry)
        text_lower = text.lower()
        if (patient_code_lower in text_lower) and re.search(r"\bmd[123]\b", text_lower, re.IGNORECASE):
            matches.append((k, entry))
    return matches

# =======================
# COMMIT (supports deletion), AUTOSAVE
# =======================
def _ensure_initial_suffix(txt: str) -> str:
    s = (txt or "").strip()
    if not s:
        return s
    if re.search(r"\binitial\s*dose\b", s, re.IGNORECASE):
        return s
    if _extract_patient_code(s) and not _is_maintenance(s):
        return s + " Initial Dose"
    return s

def _commit_and_autosave(dkey: str, widget_key: str):
    try:
        raw_val = ss.get(widget_key, "")
        old_entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
        if isinstance(old_entry, str):
            old_entry = {"text": old_entry, "cancelled": False}
            ss.entries[dkey] = old_entry
        old_text = old_entry.get("text", "")
        old_cancelled = old_entry.get("cancelled", False)

        # 🚨 Handle "delete" keyword
        if str(raw_val).strip().lower() == "delete":
            patient_code = _extract_patient_code(old_text)
            if patient_code and not _is_maintenance(old_text):
                # Delete all maintenance doses
                for md_key, _ in _find_maintenance_doses(patient_code):
                    ss.entries.pop(md_key, None)
                    md_widget_key = f"cell_widget_{md_key}"
                    if md_widget_key in ss:
                        del ss[md_widget_key]
            ss.entries.pop(dkey, None)
            ss[widget_key] = ""
            _autosave_now()
            ss["__autosave_ok__"] = True
            ss["__autosave_error__"] = ""
            return

        # 🚨 Guard: skip autosave if widget looks blank but entry still has data
        if (not raw_val or not str(raw_val).strip()) and old_text.strip():
            ss[widget_key] = old_text  # resync widget
            return

        # --- Parse input and detect cancellation ---
        text_val = str(raw_val).strip()
        cancelled = old_cancelled

        # Detect cancellation
        ends_with_cancel = text_val.lower().endswith("cancel") or text_val.lower().endswith("cancelled")
        if ends_with_cancel:
            text_val = re.sub(r"\s*[-–—]?\s*cancel(?:led)?$", "", text_val, flags=re.IGNORECASE).strip()
            cancelled = True

        val = _ensure_initial_suffix(text_val)

        # Extract patient code from new or old value
        patient_code = _extract_patient_code(val) or _extract_patient_code(old_text)

        # 🔴 If this is an initial dose being newly cancelled
        if (
            cancelled and not old_cancelled and
            patient_code and
            not _is_maintenance(old_text)
        ):
            # Mark all maintenance doses as cancelled
            for md_key, md_entry in _find_maintenance_doses(patient_code):
                # Preserve text, just update cancelled flag
                if isinstance(md_entry, dict):
                    new_md_entry = {"text": md_entry["text"], "cancelled": True}
                else:
                    new_md_entry = {"text": str(md_entry), "cancelled": True}
                ss.entries[md_key] = new_md_entry
                md_widget_key = f"cell_widget_{md_key}"
                if md_widget_key in ss:
                    ss[md_widget_key] = new_md_entry["text"]

        # 🔵 If un-cancelling an initial dose, you could re-activate MDs here (optional)
        # For now, we leave them cancelled until manually edited

        # Update current entry
        ss.entries[dkey] = {"text": val, "cancelled": cancelled}
        ss[widget_key] = val

        # Only schedule maintenance doses if it's a new/active initial dose and not cancelled
        if (
            not cancelled and
            (not old_cancelled or not old_text) and  # includes new entries
            patient_code and
            not _is_maintenance(val)
        ):
            init_dt = _parse_date_from_dkey(dkey)
            if init_dt:
                _schedule_patient_cycle(
                    patient_code=patient_code,
                    initial_dt=init_dt,
                    n_maint=3,
                    interval_weeks=6,
                    base_text=val
                )
                ss[RERUN_FLAG] = True

        _autosave_now()
        ss["__autosave_ok__"] = True
        ss["__autosave_error__"] = ""

    except Exception as e:
        ss["__autosave_ok__"] = False
        ss["__autosave_error__"] = str(e)


def _commit_all_widgets_and_autosave():
    try:
        changed = False
        prefix = "cell_widget_"
        for key in list(ss.keys()):
            if not (isinstance(key, str) and key.startswith(prefix)):
                continue
            dkey = key[len(prefix):]
            new_val_raw = ss.get(key, "")
            old_entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
            old_text = old_entry.get("text", "")
            old_cancelled = old_entry.get("cancelled", False)

            if not new_val_raw or not str(new_val_raw).strip():
                if str(new_val_raw).strip().lower() == "delete":
                    # User typed delete → clear it
                    ss.entries.pop(dkey, None)
                    ss[key] = ""
                    changed = True
                    continue

                # 🚨 Guard: only delete if widget AND entries are empty
                if dkey in ss.entries and not ss.get(key, "").strip():
                    del ss.entries[dkey]
                    ss[key] = ""
                    changed = True
                continue

            # detect cancellation keyword in raw input
            cancelled = old_cancelled
            text_val = str(new_val_raw).strip()
            if text_val.lower().endswith("cancel") or text_val.lower().endswith("cancelled"):
                text_val = re.sub(r"\s*[-–—]?\s*cancel(?:led)?$", "", text_val, flags=re.IGNORECASE).strip()
                cancelled = True

            new_val = _ensure_initial_suffix(text_val)

            if new_val != old_text or cancelled != old_cancelled:
                ss.entries[dkey] = {"text": new_val, "cancelled": cancelled}
                ss[key] = new_val
                changed = True

                code = _extract_patient_code(new_val)
                if code and not _is_maintenance(new_val):
                    init_dt = _parse_date_from_dkey(dkey)
                    if init_dt:
                        _schedule_patient_cycle(
                            patient_code=code, initial_dt=init_dt, n_maint=3, interval_weeks=6, base_text=new_val
                        )
                        ss[RERUN_FLAG] = True

        if changed:
            _autosave_now()
        ss["__autosave_ok__"] = True
        ss["__autosave_error__"] = ""
    except Exception as e:
        ss["__autosave_ok__"] = False
        ss["__autosave_error__"] = str(e)

# =======================
# BOOT LOAD + WATCHDOG
# =======================
if "__boot_loaded__" not in ss:
    ss["__boot_loaded__"] = True
    try:
        fp = _get_json_path()
        entries, meta, week_action_rows, full_data = _try_load_from(fp)
        if entries is not None:
            # entries now contain dicts: {"text": ..., "cancelled": ...}
            ss.entries = entries
            if week_action_rows is not None:
                ss.week_action_rows = week_action_rows
            if full_data and "custom_legend_entries" in full_data:
                ss.custom_legend_entries = full_data["custom_legend_entries"]
            if full_data and "suppressed_us_holidays" in full_data:
                ss.suppressed_us_holidays = full_data["suppressed_us_holidays"]                
            changed = _apply_meta_to_calendar(meta or {})
            _preload_widgets_from_entries()
            ss["__disk_mtime__"] = _stat_mtime(fp)
            if changed:
                rerun()
    except Exception as e:
        ss["__boot_error__"] = str(e)

# Persist meta if month/year changed
if "__last_meta__" not in ss:
    ss["__last_meta__"] = (ss.current_year, ss.current_month)
if ss["__last_meta__"] != (ss.current_year, ss.current_month):
    _autosave_now()
    ss["__last_meta__"] = (ss.current_year, ss.current_month)

# Check external changes
_disk_watchdog()


# =======================
# Title and NAV
# =======================
col_left, col_mid, col_right = st.columns([0.48, 5.04, 0.48])

# --- ROW 1: Legend and Title ---
with col_left:
    if st.button("☰ Legend", key="legend_toggle"):
        ss.show_legend = not ss.show_legend
        st.rerun()

with col_mid:
    st.markdown(
        "<h1 style='text-align: center; margin: 0; font-size: 40px;'>Production Schedule Dashboard</h1>",
        unsafe_allow_html=True
    )

# Leave col_right empty for symmetry

# --- ROW 2: Prev, Month, Next ---
col_nav_left, col_nav_mid, col_nav_right = st.columns([0.48, 5.04, 0.265])

with col_nav_left:
    if st.button("← Prev", key="prev"):
        _commit_all_widgets_and_autosave()
        ss.current_month -= 1
        if ss.current_month < 1:
            ss.current_month = 12
            ss.current_year -= 1
        st.rerun()

with col_nav_mid:
    st.markdown(
        f"<h2 style='text-align: center; margin: 0; font-weight: 600; color: #0047AB; font-size: 30px;'>"
        f"{calendar.month_name[ss.current_month]} {ss.current_year}</h2>",
        unsafe_allow_html=True
    )

with col_nav_right:
    if st.button("Next →", key="next"):
        _commit_all_widgets_and_autosave()
        ss.current_month += 1
        if ss.current_month > 12:
            ss.current_month = 1
            ss.current_year += 1
        st.rerun()

# Divider before calendar
st.markdown("---")

# =======================
# LEGEND SIDEBAR
# =======================
if ss.show_legend:
    with st.sidebar:
        st.title("Legends")

        if 'custom_legend_entries' not in ss:
            ss.custom_legend_entries = []

        built_in_legends = [
            {"label": "Confirmed Patient", "description": "Confirmed Patient Dose Scheduled", "color": COLOR_CONFIRMED, "builtin": True},
            {"label": "Placeholder Patient", "description": "Placeholder for Expected Patient Dose", "color": COLOR_PLACEHOLDER, "builtin": True},
            {"label": "Shutdown", "description": "Equipment or Facility Shutdown", "color": COLOR_SHUTDOWN, "builtin": True},
            {"label": "Cardinal/TPI/Niowave", "description": "Ac225 Production site activities", "color": COLOR_CARDINAL_TPI_NIOWAVE, "builtin": True},
            {"label": "BWXT Order", "description": "IN-111 Isotope", "color": COLOR_BWXT, "builtin": True},
            {"label": "AC225 Run-EVG", "description": "Scheduled production of Ac225 batches at Evergreen", "color": COLOR_AC225_RUN_EVG, "builtin": True},
            {"label": "IN111 Run-EVG", "description": "Scheduled production of In111 batches at Evergreen", "color": COLOR_IN111_RUN_EVG, "builtin": True},
            {"label": "AC225 Run-SRx", "description": "Scheduled production of Ac225 batches at Spectron Rx", "color": COLOR_AC225_RUN_SRX, "builtin": True},
            {"label": "IN111 Run-SRx", "description": "Scheduled production of In111 batches at Spectron Rx", "color": COLOR_IN111_RUN_SRX, "builtin": True},
            {"label": "NMCTG", "description": "Clinical Site Qualification Event by NMCTG", "color": COLOR_NMCTG, "builtin": True},
            {"label": "Perceptive", "description": "Clinical Site Qualification Event by Perceptive", "color": COLOR_PERCEPTIVE, "builtin": True},
            {"label": "Maintenance Dose", "description": "Maintenance Dose for Confirmed Patient", "color": COLOR_MD, "builtin": True},
            {"label": "PV SRx", "description": "Process Validation Spectron Rx", "color": COLOR_PV, "builtin": True},
            {"label": "SRx Maintenance", "description": "Spectron Rx Maintenance", "color": COLOR_SRX, "builtin": True},
        ]

        # show built-ins + customs
        for item in (built_in_legends + [
            {"label": it["label"], "description": it["description"], "color": it["color"], "builtin": False, "index": i}
            for i, it in enumerate(ss.custom_legend_entries)
        ]):
            cols = st.columns([4, 1])
            with cols[0]:
                text_color = "black" if is_light_color(item['color']) else "white"
                st.markdown(
                    f"""
                    <div style="background-color:{item['color']};padding:10px;margin:6px 0;border-radius:6px;border:1px solid #ddd;">
                        <div style="font-weight:600;color:{text_color};font-size:13px;">{item['label']}</div>
                        <div style="color:{text_color};font-size:11px;">{item['description']}</div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            with cols[1]:
                if not item.get("builtin", False):
                    if st.button("🗑️", key=f"del_custom_legend_{item['index']}"):
                        ss.custom_legend_entries.pop(item['index'])
                        rerun()

        st.markdown("### ➕ Add New Legend")
        picked_color = st.color_picker("Choose color:", "#3366cc", key="new_legend_color")
        label_input = st.text_input("Symbol/Label", placeholder="Enter Symbol/Label")
        desc_input = st.text_input("Description", placeholder="Enter Description")
        if st.button("➕ Add Legend Item", key="add_custom_legend_btn"):
            if label_input.strip():
                ss.custom_legend_entries.append({
                    "label": label_input.strip(),
                    "description": desc_input.strip() if desc_input.strip() else "No description",
                    "color": picked_color
                })
                st.success(f"Added: {label_input}")
                rerun()
            else:
                st.warning("Label is required.")

        # =======================
        # MANAGE U.S. HOLIDAYS (Dropdown Version)
        # =======================
        st.markdown("---")
        st.subheader("U.S. Federal Holidays")

        us_holidays = holidays.US(years=ss.current_year)

        # Build list of (name, date) for current year
        holiday_items = [(name, dt) for dt, name in us_holidays.items()]
        unique_holiday_names = sorted(set(name for name, dt in holiday_items))

        # Sort by date for logical order
        sorted_holiday_items = sorted(holiday_items, key=lambda x: x[1])

        # Dropdown
        selected_holiday = st.selectbox(
            "Select a Federal Holiday",
            options=[f"{name} ({dt.strftime('%b %d')})" for name, dt in sorted_holiday_items],
            index=None,
            placeholder="Choose A Holiday",
            key="select_holiday"
        )

        if selected_holiday:
            # Parse name and date from display string
            name_part = selected_holiday.rsplit(" (", 1)[0]
            date_str = selected_holiday.split(" (")[1].rstrip(")")

            # Find the actual holiday record
            selected_dt = None
            for name, dt in holiday_items:
                if name == name_part and dt.strftime("%b %d") == date_str:
                    selected_dt = dt
                    break

            if selected_dt:
                dkey = date_key(selected_dt.year, selected_dt.month, selected_dt.day, 0)
                is_suppressed = name_part in ss.suppressed_us_holidays
                status_icon = "❌ Removed" if is_suppressed else "✅ Active"
                status_color = "gray" if is_suppressed else "black"

                # Show holiday preview with formatted HTML
                st.markdown(
                    f"""
                    <small style='color: {status_color};'>
                        📅 <strong>{name_part}</strong><br>
                        Date: {selected_dt.strftime('%A, %B %d, %Y')}<br>
                        Status: {status_icon}
                    </small>
                    """,
                    unsafe_allow_html=True
                )

                # Action button
                button_label = "✅ Add Back" if is_suppressed else "🗑️ Remove"
                button_type = "primary" if is_suppressed else "secondary"

                if st.button(button_label, key=f"toggle_holiday_{name_part}", type=button_type, use_container_width=True):
                    if is_suppressed:
                        ss.suppressed_us_holidays.remove(name_part)
                    else:
                        ss.suppressed_us_holidays.append(name_part)

                    # Optional: Remove from entries if being hidden
                    if not is_suppressed:
                        if dkey in ss.entries:
                            del ss.entries[dkey]
                        widget_key = f"cell_widget_{dkey}"
                        if widget_key in ss:
                            del ss[widget_key]

                    _autosave_now()
                    st.rerun()

        # =======================
        #  (Custom Holidays)
        # =======================
        st.markdown("---")
        st.subheader("Add New Holiday")

        with st.form(key="form_add_closure"):
            closure_name = st.text_input("Closure Name", placeholder="Enter Holiday")
            closure_date = st.date_input(
                "Select Date",
                value=date.today(),
                min_value=date(2000, 1, 1),
                max_value=date(2100, 12, 31)
            )
            submit = st.form_submit_button("Add Holiday")

            if submit and closure_name.strip():
                new_closure = {
                    "name": closure_name.strip(),
                    "date": closure_date.isoformat()
                }
                if new_closure not in ss.custom_closures:
                    ss.custom_closures.append(new_closure)
                    _autosave_now()
                    st.rerun()
                else:
                    st.warning("This closure already exists.")

        # Display existing closures
        if ss.custom_closures:
            st.markdown("### Active Custom Holidays")
            for idx, closure in enumerate(ss.custom_closures):
                closure_date = date.fromisoformat(closure["date"])
                col_del, col_info = st.columns([0.8, 3])
                with col_del:
                    if st.button("🗑️", key=f"del_closure_{idx}", help="Remove Holiday"):
                        ss.custom_closures.pop(idx)
                        _autosave_now()
                        st.rerun()
                with col_info:
                    st.markdown(
                        f"<small>{closure['name']} — {closure_date.strftime('%b %d, %Y')}</small>",
                        unsafe_allow_html=True
                    )

        if st.button("✕ Close Legend", use_container_width=True, type="primary", key="close_legend"):
            ss.show_legend = False
            rerun()

# =======================
# CALENDAR GRID (Week N aligned with date row + inline +Row/-Row)
# =======================
cal_raw = calendar.monthcalendar(ss.current_year, ss.current_month)
valid_weeks = [week for week in cal_raw if any(d != 0 for d in week)]

# Extended dates (Mon..Sun with spillover for alignment)
extended_weeks = []
for week in valid_weeks:
    extended_week = []
    for i, d in enumerate(week):
        if d == 0:
            ref_day = next((x for x in week if x != 0), None)
            if ref_day is None:
                extended_week.append(None)
                continue
            ref_date = date(ss.current_year, ss.current_month, ref_day)
            monday_index = ref_date.weekday()  # 0..6
            target_date = ref_date - timedelta(days=monday_index - i)
            extended_week.append(target_date)
        else:
            try:
                extended_week.append(date(ss.current_year, ss.current_month, d))
            except Exception:
                extended_week.append(None)
    extended_weeks.append(extended_week)

# Ensure we have enough per-week rows based on saved entries
def _ensure_rows_for_current_month(valid_weeks_list):
    day_to_week = {}
    for w_idx, week in enumerate(valid_weeks_list):
        for d in week:
            if d != 0:
                day_to_week[d] = w_idx
    required_rows = {}
    for k in ss.entries.keys():
        try:
            dpart, rpart = k.split("_", 1)
            dt_ = datetime.fromisoformat(dpart).date()
            if dt_.year == ss.current_year and dt_.month == ss.current_month:
                row_i = int(rpart)
                w_idx = day_to_week.get(dt_.day, None)
                if w_idx is None:
                    continue
                needed = row_i + 1
                required_rows[w_idx] = max(required_rows.get(w_idx, 1), needed)
        except Exception:
            continue
    for w_idx in range(len(valid_weeks_list)):
        key = f"{ss.current_year}-{ss.current_month}_{w_idx}"
        current = ss.week_action_rows.get(key, 1)
        ss.week_action_rows[key] = max(current, required_rows.get(w_idx, 1))

# === ADD U.S. HOLIDAYS (Safe: Before widget sync) ===
us_holidays = holidays.US(years=ss.current_year)
for holiday_date, holiday_name in us_holidays.items():
    if holiday_date.month == ss.current_month and holiday_date.year == ss.current_year:
        # ✅ Skip if suppressed
        if holiday_name in ss.suppressed_us_holidays:
            continue
        dkey = date_key(holiday_date.year, holiday_date.month, holiday_date.day, 0)
        entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
        if isinstance(entry, str):
            # Legacy cleanup: convert old string to new format
            entry = {"text": entry.strip(), "cancelled": False}
        current_text = entry["text"].strip()

        # Only set if empty or placeholder like "Weekend"
        if not current_text or current_text.lower() == "weekend":
            entry["text"] = holiday_name
            entry["cancelled"] = False
            ss.entries[dkey] = entry

            widget_key = f"cell_widget_{dkey}"
            if widget_key not in ss:
                ss[widget_key] = holiday_name

# === ADD CUSTOM CLOSURES ===
for closure in ss.custom_closures:
    try:
        closure_date = date.fromisoformat(closure["date"])
        if closure_date.month == ss.current_month and closure_date.year == ss.current_year:
            dkey = date_key(closure_date.year, closure_date.month, closure_date.day, 0)
            entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
            if isinstance(entry, str):
                entry = {"text": entry.strip(), "cancelled": False}
            current_text = entry["text"].strip()

            if not current_text or current_text.lower() == "weekend":
                entry["text"] = closure["name"]
                entry["cancelled"] = False
                ss.entries[dkey] = entry

                widget_key = f"cell_widget_{dkey}"
                if widget_key not in ss:
                    ss[widget_key] = closure["name"]
    except Exception as e:
        continue  # Skip invalid dates

_ensure_rows_for_current_month(valid_weeks)

# === SYNC WIDGETS WITH TEXT FIELD ONLY ===
def _sync_widgets_with_entries():
    for k, v in ss.entries.items():
        if isinstance(v, dict):
            text_val = v.get("text", "")
        else:
            text_val = str(v)  # fallback for legacy
        wk = f"cell_widget_{k}"
        if wk not in ss or ss[wk] != text_val:
            ss[wk] = text_val

_sync_widgets_with_entries()

# --- Header Row: "Week No" + Day Names ---
header_cols = st.columns([0.8, 1, 1, 1, 1, 1, 1, 1])

with header_cols[0]:
    st.markdown(
        "<div style='font-size:20px; font-weight:800; text-align:center'>Week</div>",
        unsafe_allow_html=True
    )

for i, day_name in enumerate(["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"], start=1):
    with header_cols[i]:
        st.markdown(
            f"<div style='font-size:20px; font-weight:800; text-align:center'>{day_name}</div>",
            unsafe_allow_html=True
        )

st.markdown("---")

# --- Render Each Week ---
for week_idx, week_dates in enumerate(extended_weeks):
    week_key = f"{ss.current_year}-{ss.current_month}_{week_idx}"
    num_rows = ss.week_action_rows.get(week_key, 1)

    # === DATE ROW + WEEK LABEL IN THE SAME ROW (Perfect Alignment) ===
    row_cols = st.columns([0.8, 1, 1, 1, 1, 1, 1, 1])

    # Left rail: "Week N" — perfectly aligned with date cells
    with row_cols[0]:
        st.markdown(
            f"""
            <div style="
                font-size:16px;
                font-weight:600;
                text-align:center;
                line-height:40px;
                height:40px;
                display:flex;
                align-items:center;
                justify-content:center;
                margin:0;
                padding:0;
            ">
                Week {week_idx+1}
            </div>
            """,
            unsafe_allow_html=True
        )

    # Date cells (Mon-Sun) — same height and alignment
    for i, dtm in enumerate(week_dates):
        with row_cols[i + 1]:
            if dtm is None:
                st.write("")
            else:
                st.markdown(
                    f"""
                    <div style="
                        font-size:16px;
                        font-weight:600;
                        text-align:center;
                        line-height:40px;
                        height:40px;
                        display:flex;
                        align-items:center;
                        justify-content:center;
                        margin:0;
                        padding:0;
                    ">
                        {dtm.strftime('%b-%d')}
                    </div>
                    """,
                    unsafe_allow_html=True
                )

    # === EVENT ROWS ===
    for row_idx in range(num_rows):
        event_cols = st.columns([0.8, 1, 1, 1, 1, 1, 1, 1])

        # Left rail: Show +Row / -Row only in the FIRST row, spacer otherwise
        with event_cols[0]:
            if row_idx == 0:
                st.markdown('<div style="margin-top: 4px;">', unsafe_allow_html=True)
                c_add, c_del = st.columns(2)
                with c_add:
                    if st.button(
                        "➕",
                        key=f"wk_add_{ss.current_year}_{ss.current_month}_{week_idx}",
                        use_container_width=True,
                        help="Add a new row at the bottom"
                    ):
                        ss.week_action_rows[week_key] = num_rows + 1
                        _autosave_now()
                        st.rerun()

                with c_del:
                    if st.button("➖", key=f"wk_del_{ss.current_year}_{ss.current_month}_{week_idx}", use_container_width=True, help="Delete the last row (only if empty)"):
                        current_rows = ss.week_action_rows.get(week_key, 1)
                        if current_rows <= 1:
                            st.toast("This Row Cannot Be Deleted: It Is The Only Entry For This Week", icon="⚠️")
                        else:
                            bottom_row_empty = True
                            for day in valid_weeks[week_idx]:
                                if day != 0:
                                    dkey = date_key(ss.current_year, ss.current_month, day, current_rows - 1)
                                    entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
                                    text_val = entry.get("text", "").strip()
                                    # Treat "Weekend" as empty
                                    if text_val and text_val != "Weekend":
                                        bottom_row_empty = False
                                        break
                            if bottom_row_empty:
                                for day in valid_weeks[week_idx]:
                                    if day != 0:
                                        dkey = date_key(ss.current_year, ss.current_month, day, current_rows - 1)
                                        ss.entries.pop(dkey, None)
                                        ss.pop(f"cell_widget_{dkey}", None)
                                ss.week_action_rows[week_key] = current_rows - 1
                                _autosave_now()
                                st.rerun()
                            else:
                                st.toast("This Row Cannot Be Deleted: It Contains Scheduled Events. Please Remove The Events Before Deleting Row", icon="⚠️")

        # Event cells (Mon-Sun)
        for day_idx, dtm in enumerate(week_dates):
            with event_cols[day_idx + 1]:
                if dtm is None:
                    st.write("")
                else:
                    dkey = date_key(dtm.year, dtm.month, dtm.day, row_idx)
                    widget_key = f"cell_widget_{dkey}"

                    # Get current entry — always ensure dict structure
                    entry = ss.entries.get(dkey, {"text": "", "cancelled": False})
                    if isinstance(entry, str):
                        # Migrate legacy string entry
                        entry = {"text": entry.strip(), "cancelled": False}
                        ss.entries[dkey] = entry

                    text_val = entry["text"]
                    cancelled = entry["cancelled"]

                    # Handle weekend auto-fill
                    is_weekend = dtm.weekday() >= 5
                    if is_weekend:
                        if not text_val or text_val == "Weekend":
                            # Ensure weekend is set
                            entry["text"] = "Weekend"
                            entry["cancelled"] = False
                            ss.entries[dkey] = entry
                            ss[widget_key] = "Weekend"
                        # Don't allow editing if it's just "Weekend"?
                        # But let user override — so we keep input enabled

                    # Sync widget to show only text (not cancellation flag)
                    if widget_key not in ss:
                        ss[widget_key] = text_val

                    display_val = text_val

                    # Apply color using full entry dict
                    bg_color = get_color(entry)
                    text_color = "black" if is_light_color(bg_color) else "white"

                    label_str = f"cell_{dkey}"

                    # --- Render text input ---
                    st.markdown(
                        f"""
                        <style>
                        div[data-testid="stTextInput"] input[aria-label="{label_str}"] {{
                            background-color: {bg_color} !important;
                            color: {text_color} !important;
                            border: 0 !important;
                            height: 40px !important;
                            line-height: 40px !important;
                            text-align: center !important;
                            font-weight: 500 !important;
                            border-radius: 4px !important;
                            box-shadow: none !important;
                            padding: 0 8px !important;
                            margin: 0 !important;
                        }}
                        div[data-testid="stTextInput"] label {{ display: none !important; }}
                        div[data-testid="stTextInput"] > div {{ margin: 0 !important; padding: 0 !important; }}
                        </style>
                        """,
                        unsafe_allow_html=True
                    )

                    # Only update widget if it hasn't been touched
                    if ss.get(widget_key) != display_val and ss.get(widget_key) == text_val:
                        ss[widget_key] = display_val

                    st.text_input(
                        label=label_str,
                        key=widget_key,
                        label_visibility="collapsed",
                        placeholder="Add event" if not is_weekend else "",
                        on_change=lambda dk=dkey, wk=widget_key: _commit_and_autosave(dk, wk),
                    )

    # Spacing between weeks
    st.markdown('<div style="margin: 12px 0;"></div>', unsafe_allow_html=True)


# =======================
# EXPORTS
# =======================
def _safe_set_auto_size(text_frame):
    if MSO_AUTO_SIZE is None:
        return
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass  # tolerate odd python-pptx versions

def generate_pdf_calendar(year, month, entries, week_action_rows):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4, leftMargin=0.4*inch, rightMargin=0.4*inch,
        topMargin=0.5*inch, bottomMargin=0.5*inch
    )
    story = []
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=24, spaceAfter=15, alignment=1)
    month_style = ParagraphStyle('MonthStyle', parent=styles['Heading2'], fontSize=20, spaceAfter=20, alignment=1)
    story.append(Paragraph("Production Schedule Dashboard", title_style))
    story.append(Paragraph(f"{calendar.month_name[month]} {year}", month_style))

    cal_raw = calendar.monthcalendar(year, month)
    valid_weeks = [week for week in cal_raw if any(d != 0 for d in week)]
    extended_weeks = []
    for week in valid_weeks:
        ref_day = next((d for d in week if d != 0), None)
        if ref_day is None:
            extended_weeks.append([None]*7); continue
        ref_date = date(year, month, ref_day)
        week_start = ref_date - timedelta(days=ref_date.weekday())
        extended_weeks.append([week_start + timedelta(days=i) for i in range(7)])

    cell_style = ParagraphStyle('TableCell', fontSize=9, leading=10, alignment=1, wordWrap='CJK', spaceAfter=2, textColor=colors.black)
    header_style = ParagraphStyle('HeaderCell', parent=cell_style, fontSize=12, textColor=colors.whitesmoke, fontName='Helvetica-Bold')
    date_style = ParagraphStyle('DateCell', parent=cell_style, fontSize=12, textColor=colors.black, fontName='Helvetica-Bold')

    day_headers = ["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"]
    table_data = [[Paragraph(h, header_style) for h in day_headers]]
    row_heights = [0.4]

    for week_idx, week_dates in enumerate(extended_weeks):
        table_data.append([Paragraph(d.strftime('%b-%d'), date_style) if d else "" for d in week_dates])
        row_heights.append(0.35)
        num_rows = week_action_rows.get(f"{year}-{month}_{week_idx}", 1)  # Use full key
        for row_idx in range(num_rows):
            row = []
            for d in week_dates:
                if not d:
                    row.append("")
                    continue
                dk = date_key(d.year, d.month, d.day, row_idx)
                raw_entry = entries.get(dk)

                # Handle missing or empty entry
                if raw_entry is None:
                    row.append("")
                    continue

                # Normalize entry
                if isinstance(raw_entry, str):
                    text_val = raw_entry.strip()
                    cancelled = False
                else:
                    text_val = raw_entry.get("text", "").strip()
                    cancelled = raw_entry.get("cancelled", False)

                # Skip if cancelled
                if cancelled:
                    row.append("")
                    continue

                # Skip if empty
                if not text_val or text_val.lower() == "weekend":
                    row.append("")
                    continue

                # Style the paragraph
                p_style = cell_style.clone('tmp')
                color_hex = get_color(raw_entry)  # Pass full entry for color logic
                p_style.textColor = colors.black if is_light_color(color_hex) else colors.white
                row.append(Paragraph(text_val, p_style))
            table_data.append(row)
            row_heights.append(0.5)

    col_widths = [1.1*inch]*7
    table = Table(table_data, colWidths=col_widths, rowHeights=[h*inch for h in row_heights])
    table_style = TableStyle([
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
        ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
        ('GRID',(0,0),(-1,-1),1,colors.black),
        ('BACKGROUND',(0,0),(-1,0),colors.grey),
        ('BACKGROUND',(0,1),(-1,1),colors.lightgrey),
        ('LEFTPADDING',(0,0),(-1,-1),3),
        ('RIGHTPADDING',(0,0),(-1,-1),3),
        ('TOPPADDING',(0,0),(-1,-1),3),
        ('BOTTOMPADDING',(0,0),(-1,-1),3),
    ])

    current_row = 1  # Header row
    for week_idx, week_dates in enumerate(extended_weeks):
        current_row += 1  # Date row
        num_activity_rows = week_action_rows.get(f"{year}-{month}_{week_idx}", 1)
        for row_offset in range(num_activity_rows):
            for day_idx, d in enumerate(week_dates):
                if not d:
                    continue
                dk = date_key(d.year, d.month, d.day, row_offset)
                raw_entry = entries.get(dk)
                if raw_entry is None:
                    continue

                if isinstance(raw_entry, str):
                    text_val = raw_entry.strip()
                    cancelled = False
                else:
                    text_val = raw_entry.get("text", "").strip()
                    cancelled = raw_entry.get("cancelled", False)

                if cancelled or not text_val or text_val.lower() == "weekend":
                    continue

                color_hex = get_color(raw_entry)
                if color_hex == "white":
                    continue

                try:
                    h = color_hex.lstrip('#')
                    r, g, b = int(h[0:2],16)/255.0, int(h[2:4],16)/255.0, int(h[4:6],16)/255.0
                    table_style.add('BACKGROUND', (day_idx, current_row), (day_idx, current_row), colors.Color(r,g,b))
                except Exception:
                    pass
            current_row += 1

    table.setStyle(table_style)
    story.append(table)
    story.append(Spacer(1, 0.3*inch))
    doc.build(story)
    pdf_data = buffer.getvalue()
    buffer.close()
    return pdf_data

def generate_ppt_calendar(year, month, entries, week_action_rows):
    from datetime import date as dt, timedelta
    import calendar as _cal

    cal_raw = _cal.monthcalendar(year, month)
    valid_weeks = [w for w in cal_raw if any(d != 0 for d in w)]
    extended_weeks = []
    for week in valid_weeks:
        ref_day = next((d for d in week if d != 0), None)
        if ref_day is None:
            extended_weeks.append([None] * 7)
            continue
        ref_date = dt(year, month, ref_day)
        start_of_week = ref_date - timedelta(days=ref_date.weekday())
        extended_weeks.append([start_of_week + timedelta(days=i) for i in range(7)])

    prs = Presentation()
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    margin = Inches(0.4)
    prs.slide_width = int(slide_width)
    prs.slide_height = int(slide_height)

    CELL_WIDTH = (slide_width - 2 * margin) / 7
    HEADER_ROW_HEIGHT = Inches(0.3)
    DATE_ROW_HEIGHT = Inches(0.28)
    ACTIVITY_ROW_HEIGHT = Inches(0.35)
    TOP_MARGIN = Inches(1.0)
    BOTTOM_LIMIT = slide_height - margin

    def _safe_set_auto_size(text_frame):
        if MSO_AUTO_SIZE is None:
            return
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

    current_week_idx = 0
    while current_week_idx < len(extended_weeks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(margin, Inches(0.2), slide_width - 2 * margin, Inches(0.3))
        tf = title_box.text_frame
        tf.text = "Production Schedule Dashboard"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(margin, Inches(0.5), slide_width - 2 * margin, Inches(0.2))
        sf = subtitle_box.text_frame
        sf.text = f"{calendar.month_name[month]} {year}"
        sf.paragraphs[0].font.size = Pt(12)
        sf.paragraphs[0].font.bold = True
        sf.paragraphs[0].alignment = PP_ALIGN.CENTER

        y_current = TOP_MARGIN

        # Header row: Mon - Sun
        for i, day_name in enumerate(["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]):
            x = margin + i * CELL_WIDTH
            header_box = slide.shapes.add_textbox(x, y_current, CELL_WIDTH, HEADER_ROW_HEIGHT)
            hf = header_box.text_frame
            hf.text = day_name
            p = hf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = RGBColor(128, 128, 128)
            p.font.color.rgb = RGBColor(255, 255, 255)

        y_current += HEADER_ROW_HEIGHT

        # Add weeks until full slide is filled
        while current_week_idx < len(extended_weeks):
            week_dates = extended_weeks[current_week_idx]
            week_key = f"{year}-{month}_{current_week_idx}"
            num_activity_rows = week_action_rows.get(week_key, 1)
            week_height = DATE_ROW_HEIGHT + (num_activity_rows * ACTIVITY_ROW_HEIGHT)
            if y_current + week_height > BOTTOM_LIMIT:
                break

            # Date row
            for day_idx, dt_obj in enumerate(week_dates):
                if dt_obj is None:
                    continue
                x = margin + day_idx * CELL_WIDTH
                date_box = slide.shapes.add_textbox(x, y_current, CELL_WIDTH, DATE_ROW_HEIGHT)
                df = date_box.text_frame
                df.text = dt_obj.strftime("%b-%d")
                p = df.paragraphs[0]
                p.font.size = Pt(8)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                date_box.fill.solid()
                date_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
                p.font.color.rgb = RGBColor(0, 0, 0)

            y_current += DATE_ROW_HEIGHT

            # Activity rows
            for row_idx in range(num_activity_rows):
                for day_idx, dt_obj in enumerate(week_dates):
                    if dt_obj is None:
                        continue
                    x = margin + day_idx * CELL_WIDTH
                    y = y_current
                    dkey = date_key(dt_obj.year, dt_obj.month, dt_obj.day, row_idx)
                    raw_entry = entries.get(dkey)

                    # Normalize entry
                    if raw_entry is None:
                        continue

                    if isinstance(raw_entry, str):
                        text_val = raw_entry.strip()
                        cancelled = False
                    else:
                        text_val = raw_entry.get("text", "").strip()
                        cancelled = raw_entry.get("cancelled", False)

                    # Skip if cancelled or empty or just "Weekend"
                    if cancelled or not text_val or text_val.lower() == "weekend":
                        continue

                    # Create shape and set text
                    activity_box = slide.shapes.add_textbox(x, y, CELL_WIDTH, ACTIVITY_ROW_HEIGHT)
                    af = activity_box.text_frame
                    af.text = text_val
                    af.word_wrap = True
                    _safe_set_auto_size(af)
                    p = af.paragraphs[0]
                    p.font.size = Pt(8)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    # Apply color
                    color_hex = get_color(raw_entry)  # Pass full entry for correct color
                    if color_hex != "white":
                        try:
                            h = color_hex.lstrip('#')
                            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                            activity_box.fill.solid()
                            activity_box.fill.fore_color.rgb = RGBColor(r, g, b)
                            p.font.color.rgb = RGBColor(0, 0, 0) if is_light_color(color_hex) else RGBColor(255, 255, 255)
                        except Exception:
                            pass

                y_current += ACTIVITY_ROW_HEIGHT

            current_week_idx += 1

    buffer = io.BytesIO()
    prs.save(buffer)
    ppt_data = buffer.getvalue()
    buffer.close()
    return ppt_data

def generate_excel_calendar(year, month, entries, week_action_rows):
    # Require openpyxl at runtime; show a friendly error if missing
    try:
        from openpyxl.styles import PatternFill, Font, Alignment
        from openpyxl.utils import get_column_letter
        from openpyxl import Workbook
    except Exception as e:
        raise RuntimeError("Excel export requires 'openpyxl'. Install via: pip install openpyxl") from e

    import calendar as _cal
    from datetime import date as dt, timedelta

    cal_raw = _cal.monthcalendar(year, month)
    valid_weeks = [w for w in cal_raw if any(d != 0 for d in w)]
    extended_weeks = []
    for week in valid_weeks:
        ref_day = next((d for d in week if d != 0), None)
        if ref_day is None:
            extended_weeks.append([None]*7)
            continue
        ref_date = dt(year, month, ref_day)
        start_of_week = ref_date - timedelta(days=ref_date.weekday())
        extended_weeks.append([start_of_week + timedelta(days=i) for i in range(7)])

    wb = Workbook()
    ws = wb.active
    ws.title = f"{calendar.month_name[month]} {year}"

    for col_idx in range(1, 8):
        ws.column_dimensions[get_column_letter(col_idx)].width = 15

    headers = ["Mon","Tue","Wed","Thu","Fri","Sat","Sun"]
    for col_idx, day_name in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.value = day_name
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="808080", end_color="808080", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")

    def hex_to_xlsx(color_hex):
        return color_hex.lstrip('#').upper() if color_hex != "white" else "FFFFFF"

    current_row = 2
    for week_idx, week_dates in enumerate(extended_weeks):
        # Date row
        for col_idx, dt_obj in enumerate(week_dates, 1):
            if dt_obj is None:
                continue
            cell = ws.cell(row=current_row, column=col_idx)
            cell.value = dt_obj.strftime("%b-%d")
            cell.font = Font(bold=True, size=11, color="000000")
            cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")

        num_rows = week_action_rows.get(f"{year}-{month}_{week_idx}", 1)  # Use full key
        for row_offset in range(num_rows):
            current_row += 1
            for col_idx, dt_obj in enumerate(week_dates, 1):
                if dt_obj is None:
                    continue
                dkey = date_key(dt_obj.year, dt_obj.month, dt_obj.day, row_offset)
                raw_entry = entries.get(dkey)

                # Normalize entry
                if raw_entry is None:
                    cell_value = ""
                else:
                    if isinstance(raw_entry, str):
                        text_val = raw_entry.strip()
                        cancelled = False
                    else:
                        text_val = raw_entry.get("text", "").strip()
                        cancelled = raw_entry.get("cancelled", False)

                    # Skip if cancelled or empty or just "Weekend"
                    if cancelled or not text_val or text_val.lower() == "weekend":
                        cell_value = ""
                    else:
                        cell_value = text_val

                cell = ws.cell(row=current_row, column=col_idx)
                cell.value = cell_value
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                cell.font = Font(size=10, bold=True)

                # Apply color only if there's content
                if cell_value:
                    color_hex = get_color(raw_entry)  # Use full entry for color logic
                    if color_hex != "white":
                        bg = hex_to_xlsx(color_hex)
                        text_color = "000000" if is_light_color(color_hex) else "FFFFFF"
                        cell.fill = PatternFill(start_color=bg, end_color=bg, fill_type="solid")
                        cell.font = Font(size=10, bold=True, color=text_color)
                    else:
                        cell.fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
                else:
                    cell.fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

        current_row += 1

    # Auto-fit column widths
    for col in ws.columns:
        max_len = 0
        col_letter = get_column_letter(col[0].column)
        for c in col:
            try:
                val_len = len(str(c.value)) if c.value else 0
                max_len = max(max_len, val_len)
            except Exception:
                pass
        ws.column_dimensions[col_letter].width = min(max_len + 2, 22)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()

# =======================
# EXPORT SECTION
# =======================
st.markdown("---")
st.markdown("### 📤 Export Production Schedule Dashboard")
#st.markdown("Select Your Export Format:")
month_week_rows = {i: ss.week_action_rows.get(f"{ss.current_year}-{ss.current_month}_{i}", 1) for i in range(len(valid_weeks))}
with st.container():
    if "export_data" not in ss:
        ss.export_data = {}
    current_month_key = f"{ss.current_year}-{ss.current_month}"
    if ss.export_data.get("_month") != current_month_key:
        ss.export_data = {"_month": current_month_key}

    formats = [
        {"name": "PowerPoint", "icon": "", "key": "ppt", "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "ext": "pptx", "generator": generate_ppt_calendar, "label": "Generate PowerPoint"},
        {"name": "PDF", "icon": "", "key": "pdf", "mime": "application/pdf", "ext": "pdf", "generator": generate_pdf_calendar, "label": "Generate PDF"},
        {"name": "Excel", "icon": "", "key": "excel", "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "ext": "xlsx", "generator": generate_excel_calendar, "label": "Generate Excel"},
    ]

    formats_dict = {fmt["name"]: fmt for fmt in formats}

    format_options = [fmt["name"] for fmt in formats]

    selected_format = st.selectbox("Select Your Export Format :", format_options)

    fmt = formats_dict[selected_format]

    if st.button("Prepare Your Export", key="generate_button", type="secondary"):
        with st.spinner(f"🔧 Generating {selected_format}..."):
            try:
                data = fmt["generator"](ss.current_year, ss.current_month, ss.entries, month_week_rows)
                ss.export_data[fmt["key"]] = data
                st.success(f"✅ {selected_format} Ready for Download!")
            except Exception as e:
                st.error(f"❌ {selected_format} Error: {str(e)}")

    st.markdown("### Available Downloads")
    for fmt in formats:
        if fmt["key"] in ss.export_data:
            st.download_button(
                label=f"{fmt['icon']} Download {fmt['name']}",
                data=ss.export_data[fmt["key"]],
                file_name=f"production_schedule_{ss.current_year}_{ss.current_month:02d}.{fmt['ext']}",
                mime=fmt["mime"],
                key=f"dl_{fmt['key']}",
                type="primary"
            )

# ------------------------
# MANUAL SAVE BUTTON
# ------------------------
st.markdown("---")
st.markdown("### Manually Save Your Production Schedule Dashboard")
if st.button("💾 Save Now", key="manual_save_button", help="Save changes now"):
    _autosave_now()
    # Only show toast — no banners, no noise

# =======================
# RELOAD PREVIOUS
# =======================
st.markdown("---")
st.markdown("### 🔄 Reload Previous Production Schedule")
prefill_dir = str(load_latest_dir())
raw_dir_input = st.text_input(
    "Absolute directory (auto-loads if a saved file exists):",
    value=prefill_dir,
    key="dir_input",
    help="Example (macOS): /Users/you/Schedules | Example (Windows): D:\\Schedules"
)
entered_dir = _sanitize_dir(raw_dir_input)
if str(entered_dir) and str(entered_dir) != prefill_dir:
    save_latest_dir(str(entered_dir))
    new_fp = entered_dir / FILENAME
    try:
        if new_fp.exists():
            data = json.loads(new_fp.read_text(encoding="utf-8"))
            if isinstance(data, dict) and "entries" in data:
                ss["__pending_entries__"] = data.get("entries", {}) or {}
                ss["__pending_meta__"] = data.get("meta") or {}
                ss["__pending_week_action_rows__"] = data.get("week_action_rows", {}) or {}
                st.success(f"Loaded schedule from: {new_fp}")
            else:
                st.info(f"No saved schedule found at: {new_fp}")
        else:
            st.info(f"No saved schedule found at: {new_fp}")
    except Exception as e:
        st.error(f"Failed to read: {e}")
    rerun()

latest_dir = load_latest_dir()
file_path = latest_dir / FILENAME
st.markdown(
    f"<pre style='text-align:left; white-space:pre-wrap; margin:0'>Working file:\n{file_path}</pre>",
    unsafe_allow_html=True
)

# === Enhanced Save Status ===
fp = _get_json_path()
mtime = _stat_mtime(fp)

if mtime and ss.get("__disk_mtime__") == mtime:
    if ss["__autosave_ok__"]:
        st.caption("✅ All changes saved")
    else:
        st.caption("🟡 Last save had an issue")
elif mtime:
    st.caption("🔁 Changed since load — saving...")
else:
    st.caption("🆕 No file on disk yet")

if ss["__autosave_error__"]:
    st.error(f"❌ Save failed: {ss['__autosave_error__']}")

if ss.get("__boot_error__"):
    st.error(f"⚠️ Load error: {ss['__boot_error__']}")

# Manual Reload button
if st.button("↻ Reload from disk", key="reload_disk"):
    if _load_from_disk_into_state():
        ss["__disk_mtime__"] = _stat_mtime(_get_json_path())
        ss[RERUN_FLAG] = True
    else:
        st.info("No schedule file found to reload.")

# FINAL SAFE RERUN
if st.session_state.get(RERUN_FLAG):
    st.session_state[RERUN_FLAG] = False
    st.rerun()

